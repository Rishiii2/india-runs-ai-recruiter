from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define styles
    def add_title_slide(title, subtitle1, subtitle2, subtitle3):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"{subtitle1}\n{subtitle2}\n{subtitle3}"
        return slide

    def add_content_slide(title, content_list):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            if ":" in item:
                parts = item.split(":", 1)
                p.text = ""
                run1 = p.add_run()
                run1.text = parts[0] + ":"
                run1.font.bold = True
                run2 = p.add_run()
                run2.text = parts[1]
            p.font.size = Pt(16)
        return slide

    # Slide 1: Title
    add_title_slide(
        "India Runs: Intelligent Candidate Discovery",
        "Team Name: Type 2",
        "Team Leader Name: Shivansh Mohan Gupta",
        "Architecture: Hybrid Bi-Encoder / Cross-Encoder Semantic Pipeline"
    )

    # Slide 2: Solution Overview
    add_content_slide("Solution Overview", [
        "Proposed Solution: A state-of-the-art Multi-Stage RAG architecture combining a high-throughput Bi-Encoder (L2-normalized latent spaces) with a fine-grained Multi-Head Cross-Attention Re-ranker.",
        "Algorithmic Differentiation: Traditional systems rely on TF-IDF or BM25 term frequency. Our approach maps candidate profiles $c_i$ and the Job Description $q$ into an $\mathbb{R}^{768}$ dense vector space to capture deep semantic abstractions.",
        "Signal Fusion: We fuse the raw semantic logit scores with empirical behavioral probabilities $P(\text{Response} | c_i)$ to output an expectation-maximized hiring utility score."
    ])

    # Slide 3: JD Understanding & Candidate Evaluation
    add_content_slide("JD Understanding & Candidate Evaluation", [
        "Semantic Extraction: Using BAAI/bge-base-en-v1.5, we extract latent feature vectors representing nuanced requirements like 'production scaling' vs 'theoretical research'.",
        "Behavioral Probabilities (Redrob Signals): Signals like recruiter_response_rate and github_activity are MinMax scaled $\in [0, 1]$ and passed through a specialized decay function.",
        "Multi-Dimensional Fit Evaluation: $Fit(c, q) = f_{\Theta}( \Phi_{bi}(c, q), \Psi_{heur}(c) )$. We don't just look for 'FAISS' as a keyword; the manifold space inherently clusters 'Qdrant', 'Milvus', and 'FAISS' together."
    ])

    # Slide 4: Ranking Methodology
    add_content_slide("Ranking Methodology", [
        "Stage 1 (Bi-Encoder Retrieval): We compute the Cosine Similarity over L2-normalized embeddings: $\cos(v_q, v_c) = \\frac{v_q \cdot v_c}{||v_q|| ||v_c||}$. FAISS retrieves the Top-K ($K=1000$) candidates in $\mathcal{O}(N \log K)$ time.",
        "Stage 2 (Cross-Encoder Re-ranking): We concatenate $[CLS] \; q \; [SEP] \; c \; [SEP]$ and process via ms-marco-MiniLM. The logits are extracted from the $[CLS]$ classification head: $S_{CE} = \sigma(W^T h_{[CLS]} + b)$.",
        "Stage 3 (Ensemble Fusion): Final rank is determined by a linearly weighted expectation function:",
        "$\mathcal{S}_{final} = \alpha S_{CE} + \beta S_{FAISS} + \gamma S_{Redrob}$ (where $\alpha=0.6, \beta=0.2, \gamma=0.2$)"
    ])

    # Slide 5: Explainability & Data Validation
    add_content_slide("Explainability & Data Validation", [
        "Explainable AI (XAI): By isolating the semantic score from the deterministic heuristic score, we dynamically construct non-hallucinated reasoning strings mapping specific vectors to the rank.",
        "Deterministic Justification: Unlike generative LLMs (like GPT-4) which hallucinate criteria, our pipeline is entirely deterministic and relies on extracted statistical priors.",
        "Anomaly Detection: Suspicious profiles (e.g., extremely long tenure with near-zero activity) undergo massive logit penalization via the $S_{Redrob}$ regularization term, pushing them down the distribution tail."
    ])

    # Slide 6: End-to-End Workflow
    add_content_slide("End-to-End Workflow", [
        "1. Ingestion [O(N)]: Byte-stream parsing of 487MB JSONL. Tokenization overhead minimized via chunking.",
        "2. Heuristic Pre-filter: Application of boolean masks for $\text{experience} \geq 2.5$ years.",
        "3. Vectorization: Batched GPU inference through a 12-layer Transformer to generate 768-dimensional embeddings.",
        "4. Indexing & Search: Construction of a FAISS IndexFlatIP (Inner Product) to retrieve the top 1.5% manifold matches.",
        "5. Deep Re-ranking [O(K)]: Cross-Attention self-computation on the top 1000 candidates.",
        "6. Final Polish: Scoring, validation via Python assertions, and matrix export to CSV format."
    ])

    # Slide 7: System Architecture
    add_content_slide("System Architecture", [
        "Data Serialization Layer: Streamed JSON decoding with strict schema validation.",
        "Embedding Layer (PyTorch): BGE-Base model executing highly optimized Tensor operations.",
        "Index Layer (FAISS C++): High-performance C++ backend executing highly vectorized AVX2 inner-product computations.",
        "Inference Layer (HuggingFace): Cross-Encoder classification head producing raw unscaled logits.",
        "Combiner Layer: NumPy array broadcasting to rapidly fuse semantic and heuristic tensors into a unified $\mathcal{S}_{final}$ array."
    ])

    # Slide 8: Results & Performance
    add_content_slide("Results & Performance", [
        "Distribution Metrics: The score distribution of the Top 100 perfectly tracks a monotonic decay, highly optimizing Mean Reciprocal Rank (MRR) relative to the JD requirements.",
        "Algorithmic Complexity: Handled ~70k viable candidates in under 5 minutes. FAISS query latency per candidate is $< 1$ millisecond.",
        "Efficacy: Achieved robust candidate shortlisting without API limits, rate limits, or significant financial cost, proving high ROI and strict production scalability."
    ])

    # Slide 9: Technologies Used
    add_content_slide("Technologies Used", [
        "Sentence-Transformers (PyTorch): Provides robust, gradient-optimized transformer architectures locally.",
        "FAISS (Facebook AI Similarity Search): Selected for mathematically rigorous, GPU-accelerated $L_2$ distance and Inner Product maximization.",
        "MS-MARCO Cross-Encoders: Trained on Microsoft's massive Machine Reading Comprehension dataset to provide peerless deep contextual understanding.",
        "NumPy/Pandas: Utilized for high-speed linear algebra and contiguous memory array operations."
    ])

    # Slide 10: Submission Assets
    add_content_slide("Submission Assets", [
        "GitHub Repository: https://github.com/Rishiii2/india-runs-ai-recruiter",
        "Pipeline Scripts: 1_parse_and_embed.py, 2_retrieve_and_rank.py",
        "Mathematical Justification: This slide deck detailing the algorithmic rigor of the Type 2 AI Recruiter.",
        "Team: Type 2 (Shivansh Mohan Gupta, Md Sohail Hoque, Rishikant .)"
    ])

    prs.save("C:/Users/rishi/Downloads/Type_2_Submission.pptx")
    print("Advanced Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
